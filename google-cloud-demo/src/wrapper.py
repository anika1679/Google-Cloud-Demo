from google import genai
import os
from . import config
import json
import re
import sys
from pptx import Presentation
from pptx.util import Inches, Pt
import pandas as pd

# Helper Method to create Json responses from the AI API response
def parseResponseData(response):
    # Parses the information from the response, exlcuding extra text that isn't part of the json format
    json_text_match = re.search(r"\{.*\}", response, re.DOTALL)
    if json_text_match:
        json_text = json_text_match.group(0)
        # Checks to make sure it is a valid json
        try:
            data = json.loads(json_text)
            print("Parsed Successfully")
            return data
        except json.JSONDecodeError as e:
            print("Warning: Not valid JSON", e) 
            return None
        
# Method to create the Slide deck for the company analysis
def addSlide(presentation, title, content):
    # Creates the initial content slide, with no information except the title.
    content_slide = presentation.slide_layouts[1]
    slide = presentation.slides.add_slide(content_slide)
    title_shape = slide.shapes.title
    content_placeholder = slide.placeholders[1]
    title_shape.text = title
    tf = content_placeholder.text_frame
    tf.clear()
    # Adds content to the slide based on the keys from the json.
    if isinstance(content, dict):
        if "overview" in content:
            p = tf.add_paragraph()
            p.text = "Overview: " + content["overview"]
            p.font.size = Pt(14)

        if "key_stats" in content:
            tf.add_paragraph().text = "Key Stats:"
            for stat in content["key_stats"]:
                p = tf.add_paragraph()
                p.text = f"- {stat['metric']}: {stat['value']}"
                p.level = 1
        if "positive_implications" in content or "negative_implications_of_inaction" in content:
            if "positive_implications" in content:
                tf.add_paragraph().text = "Positive Implications:"
                for item in content["positive_implications"]:
                    p = tf.add_paragraph()
                    p.text = f"- {item['point']}: {item['details']}"
                    p.font.size = Pt(14)

            if "negative_implications_of_inaction" in content:
                tf.add_paragraph().text = "Negative Implications of Inaction:"
                for item in content["negative_implications_of_inaction"]:
                    p = tf.add_paragraph()
                    p.text = f"- {item['point']}: {item['details']}"
                    p.font.size = Pt(14)

    elif isinstance(content, list):
        for item in content:
            if isinstance(item, dict):
                point = item.get("point") or item.get("recommendation")
                details = item.get("details") or item.get("action")
                p = tf.add_paragraph()
                p.text = f"{point}: {details}"
                p.font.size = Pt(14)

    return slide

# Initializes the ai client being used.
client = genai.Client(vertexai=True, project=config.PROJECT_ID, location="global")
# Gets the information of the company and its compeitors from the user
company_name = input("Name of a company: ")

# Ask user for output format
print("\nChoose output format:")
print("1. PowerPoint")
print("2. CSV Table")
print("3. JSON\n")
output_choice = input("Enter your choice (1-3): ").strip()

# Prompt to be used by the AI client
prompt = f"""Think of yourself as a business analyst for {company_name}. Give:
- Overview and key stats
- SWOT analysis, listing their strenghts, weaknesses, opportunities, and risks (each on a different slide)
- Strategic recommendations for {company_name}
- Implications
Return your answer strictly in JSON using this format:
{{
  "reportTitle": "string (title of report, e.g., 'Google Business Analysis')",
  "analystPersona": "string (persona, e.g., 'Prepared by AI Business Analyst')",
  "slides": [
    {{
      "slide_title": "string (e.g., 'Overview & Key Stats')",
      "content": {{
        "overview": "string (short paragraph)",
        "key_stats": [
            {{ "metric": "string", "value": "string" }}
        ]
      }}
    }},
    {{
      "slide_title": "SWOT Analysis: Strengths",
      "content": [
        {{ "point": "string", "details": "string" }}
      ]
    }},
    {{
      "slide_title": "SWOT Analysis: Weaknesses",
      "content": [
        {{ "point": "string", "details": "string" }}
      ]
    }},
    {{
      "slide_title": "SWOT Analysis: Opportunities",
      "content": [
        {{ "point": "string", "details": "string" }}
      ]
    }},
    {{
      "slide_title": "SWOT Analysis: Risks",
      "content": [
        {{ "point": "string", "details": "string" }}
      ]
    }},
    {{
      "slide_title": "Strategic Recommendations",
      "content": [
        {{ "recommendation": "string", "action": "string" }}
      ]
    }},
    {{
      "slide_title": "Implications of Strategic Recommendations",
      "content": {{
        "positive_implications": [
          {{ "point": "string", "details": "string" }}
        ],
        "negative_implications_of_inaction": [
          {{ "point": "string", "details": "string" }}
        ]
      }}
    }}
  ]
}}
Each topic must be a new slide section. Make content concise and actionable. 
Each slide should have a clear title. 
Use bullet-like structure in "points" or "recommendations".
Do not include extra commentary outside the JSON.
"""

# Gets the response from the client
response = client.models.generate_content(
    model="gemini-2.5-flash", contents=prompt)
response_text = response.text
# Creates the data from AI response into a json that can be parsed
analysis_json = parseResponseData(response_text)
# Checks if the json was returned correctly
if not analysis_json:
    print("Error: could not parse data")
    exit()
# Executes the set of actions corressponding to the input the user gave for the output format.
if output_choice == "1":
    # Creates the presentation file
    presentation = Presentation()
    # Adds the title slide to the presentation based on the company's name
    title_slide = presentation.slide_layouts[0]
    slide = presentation.slides.add_slide(title_slide)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = analysis_json["reportTitle"]
    subtitle.text = analysis_json["analystPersona"]
    # Adds the slides for each slide section in the JSON
    for slide_data in analysis_json["slides"]:
        addSlide(presentation, slide_data["slide_title"], slide_data["content"])
    # Creates the folder to hold the powerpoints the user created
    os.makedirs("powerpoints", exist_ok=True)
    # Names the slide deck file
    presentation_name = f"{company_name}_analysis.pptx"
    # Puts the presentation in the associated folder
    file_path = os.path.join("powerpoints", presentation_name)
    # Saves presentation
    presentation.save(file_path)
    print(f"PowerPoint saved: {presentation_name}")

    if sys.platform == "darwin":
        os.system(f'open "{file_path}"')
    elif sys.platform == "win32":
        os.startfile(file_path)
    else:
        os.system(f'xdg-open "{file_path}"')

elif output_choice == "2":
    # Creates the rows for the dataframe
    rows = []
    for slide_data in analysis_json["slides"]:
        slide_title = slide_data["slide_title"]
        content = slide_data["content"]
        # Adds the content to the rows based on the keys from the json
        if isinstance(content, dict):
            if "overview" in content:
                rows.append({"Section": slide_title, "Type": "Overview", "Content": content["overview"]})
            if "key_stats" in content:
                for stat in content["key_stats"]:
                    rows.append({"Section": slide_title, "Type": "Stat", "Content": f"{stat['metric']}: {stat['value']}"})
            if "positive_implications" in content:
                for item in content["positive_implications"]:
                    rows.append({"Section": slide_title, "Type": "Positive", "Content": f"{item['point']}: {item['details']}"})
            if "negative_implications_of_inaction" in content:
                for item in content["negative_implications_of_inaction"]:
                    rows.append({"Section": slide_title, "Type": "Negative", "Content": f"{item['point']}: {item['details']}"})
        elif isinstance(content, list):
            for item in content:
                point = item.get("point") or item.get("recommendation")
                details = item.get("details") or item.get("action")
                rows.append({"Section": slide_title, "Type": "Point", "Content": f"{point}: {details}"})
    # Creates the data frame
    df = pd.DataFrame(rows)
    os.makedirs("tables", exist_ok=True)
    csv_name = f"{company_name}_analysis.csv"
    csv_path = os.path.join("tables", csv_name)
    # Makes the data frame into a csv file
    df.to_csv(csv_path, index=False)
    print(f"CSV saved: {csv_name}")

elif output_choice == "3":
    # Makes a folder for the json outputs.
    os.makedirs("jsons", exist_ok=True)
    # Names the json file
    json_name = f"{company_name}_analysis.json"
    json_path = os.path.join("jsons", json_name)
    with open(json_path, 'w') as f:
        json.dump(analysis_json, f, indent=2)
    print(f"JSON saved: {json_name}")

else:
    print(f"Invalid choice '{output_choice}'. Please enter 1, 2, or 3.")

